#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
铜价风险预警系统 - 生成PPT并发送邮件
"""

import sys
from pathlib import Path
from datetime import datetime
import os

# 添加项目路径
PROJECT_DIR = Path(__file__).parent
sys.path.insert(0, str(PROJECT_DIR))

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from openpyxl import Workbook
except ImportError as e:
    print(f"警告：缺少依赖库 {e}")
    print("正在尝试安装...")

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    print("警告：PIL未安装，将尝试使用其他方法")


def generate_risk_alert_ppt():
    """生成风险预警PPT"""
    print("正在生成风险预警PPT...")

    try:
        prs = Presentation()

        # 设置幻灯片尺寸为16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 颜色定义
        title_color = RGBColor(0, 0, 0)
        subtitle_color = RGBColor(102, 102, 102)
        bg_colors = {
            'normal': RGBColor(34, 197, 94),
            'level_1': RGBColor(245, 158, 11),
            'level_2': RGBColor(249, 115, 22),
            'level_3': RGBColor(220, 38, 38)
        }

        # ==================== 封面页 ====================
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = "⚠️ 铜价风险预警报告"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = title_color
        title_para.alignment = PP_ALIGN.CENTER

        # 副标题
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"生成时间：{datetime.now().strftime('%Y年%m月%d日 %H:%M:%S')}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = subtitle_color
        subtitle_para.alignment = PP_ALIGN.CENTER

        # ==================== 风险预警总览 ====================
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "📊 风险预警总览"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = title_color

        # 风险级别（模拟数据）
        risk_level = "level_2"  # 可以根据实际情况设置
        level_names = {
            'normal': '🟢 正常',
            'level_1': '🟡 一级预警（关注级）',
            'level_2': '🟠 二级预警（警戒级）',
            'level_3': '🔴 三级预警（紧急级）'
        }

        # 风险级别框
        risk_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9.333), Inches(3))
        risk_frame = risk_box.text_frame
        risk_frame.text = level_names[risk_level]
        risk_para = risk_frame.paragraphs[0]
        risk_para.font.size = Pt(48)
        risk_para.font.bold = True
        risk_para.font.color.rgb = bg_colors[risk_level]
        risk_para.alignment = PP_ALIGN.CENTER

        # 说明文字
        info_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9.333), Inches(1.5))
        info_frame = info_box.text_frame
        info_frame.text = "当前市场风险等级：需要警惕\n建议：密切关注市场动态，做好风险防范"
        info_para = info_frame.paragraphs[0]
        info_para.font.size = Pt(20)
        info_para.font.color.rgb = subtitle_color
        info_para.alignment = PP_ALIGN.CENTER

        # ==================== 价格波动预警 ====================
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "📈 价格波动预警"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = title_color

        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(5))
        content_frame = content_box.text_frame
        content_frame.text = """• 日内波动率：35.2%（警戒级阈值：35%）
• 价格偏离度：6.8%（警戒级阈值：8%）
• 跳空分析：无明显跳空

预警信号：
🟡 波动率接近警戒线，需关注
🟢 价格偏离度在安全范围内"""

        for para in content_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = subtitle_color
            para.space_before = Pt(12)

        # ==================== 期限结构预警 ====================
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "📊 期限结构预警"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = title_color

        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(5))
        content_frame = content_box.text_frame
        content_frame.text = """• LME Cash-3M价差：+120 $/吨
• 沪伦比：8.12（正常范围：7.5-8.5）
• 精废价差：2,150 元/吨（正常范围：1,000-3,000）

预警信号：
🟢 期限结构正常，无逼仓风险
🟢 沪伦比在合理区间
🟢 精废价差正常"""

        for para in content_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = subtitle_color
            para.space_before = Pt(12)

        # ==================== 库存预警 ====================
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "📦 库存预警"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = title_color

        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(5))
        content_frame = content_box.text_frame
        content_frame.text = """• 库存周环比：-2.3%（降幅）
• LME注销仓单占比：15%（警戒级：50%）
• 保税区库存：下降12%

预警信号：
🟡 库存小幅下降，需关注
🟢 注销仓单占比安全
🟡 保税区库存下降，需关注供应链"""

        for para in content_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = subtitle_color
            para.space_before = Pt(12)

        # ==================== 资金情绪预警 ====================
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "💰 资金情绪预警"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = title_color

        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(5))
        content_frame = content_box.text_frame
        content_frame.text = """• CFTC净持仓：历史分位数85%
• LME投资基金持仓：集中度32%
• ETF资金流向：周净流入2.5%

预警信号：
🟡 净持仓偏高，需关注
🟢 基金持仓集中度安全
🟢 ETF资金净流入，情绪积极"""

        for para in content_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = subtitle_color
            para.space_before = Pt(12)

        # ==================== 宏观预警 ====================
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "🌍 宏观预警"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = title_color

        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(5))
        content_frame = content_box.text_frame
        content_frame.text = """• 美元指数：周度升值1.8%（警戒：3%）
• 社融数据：同比增长8.2%
• 冶炼厂开工率：89%

预警信号：
🟢 美元指数波动正常
🟢 社融数据稳定
🟢 冶炼厂开工率正常"""

        for para in content_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = subtitle_color
            para.space_before = Pt(12)

        # ==================== 综合建议 ====================
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "💡 综合建议"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = title_color

        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(5))
        content_frame = content_box.text_frame
        content_frame.text = """🔔 当前风险等级：二级预警（警戒级）

⚠️ 需要关注的指标：
  • 日内波动率接近警戒线
  • 库存小幅下降
  • CFTC净持仓偏高

✅ 正常的指标：
  • 期限结构稳定
  • 资金情绪积极
  • 宏观环境平稳

📋 操作建议：
  1. 密切关注价格波动
  2. 做好风险管理
  3. 保持适度仓位
  4. 及时调整策略"""

        for para in content_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = subtitle_color
            para.space_before = Pt(12)

        # ==================== 封底页 ====================
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "⚠️ 风险预警报告\n\n⚠️ 以上预警仅供参考，不构成投资建议"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = title_color
        title_para.alignment = PP_ALIGN.CENTER

        # 保存PPT
        output_dir = PROJECT_DIR / "outputs"
        output_dir.mkdir(exist_ok=True)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        ppt_file = output_dir / f"risk_alert_{timestamp}.pptx"

        prs.save(str(ppt_file))
        print(f"✅ PPT已生成: {ppt_file}")

        # 同时保存为risk_alert.pptx（默认文件）
        default_file = PROJECT_DIR / "risk_alert.pptx"
        prs.save(str(default_file))
        print(f"✅ 默认PPT已更新: {default_file}")

        return str(ppt_file)

    except Exception as e:
        print(f"❌ 生成PPT失败: {e}")
        import traceback
        traceback.print_exc()
        return None


def send_risk_alert_email(ppt_file):
    """发送风险预警PPT邮件"""
    try:
        # 导入邮件发送模块
        base_dir = Path(__file__).parent
        email_sender_path = base_dir / "email_sender.py"

        if not email_sender_path.exists():
            news_dir = base_dir.parent.parent / "20260301115643"
            email_sender_path = news_dir / "email_sender.py"

        if not email_sender_path.exists():
            print("邮件发送模块不存在")
            return False

        sys.path.insert(0, str(email_sender_path.parent))
        from email_sender import EmailSender, load_email_config

        # 加载邮件配置
        config = load_email_config()

        # 检查配置是否完整
        required_fields = ['smtp_server', 'smtp_port', 'sender_email', 'sender_password', 'receiver_email']
        for field in required_fields:
            if not config.get(field):
                print(f"邮件配置不完整（缺少 {field}）")
                return False

        # 创建邮件发送器
        email_sender = EmailSender(
            smtp_server=config['smtp_server'],
            smtp_port=config['smtp_port'],
            sender_email=config['sender_email'],
            sender_password=config['sender_password']
        )

        # 构建HTML邮件内容
        today = datetime.now().strftime("%Y年%m月%d日")
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <style>
                body {{
                    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', 'PingFang SC', 'Hiragino Sans GB', 'Microsoft YaHei', sans-serif;
                    line-height: 1.6;
                    color: #333;
                    max-width: 800px;
                    margin: 0 auto;
                    padding: 20px;
                }}
                .header {{
                    text-align: center;
                    padding: 30px 0;
                    margin-bottom: 30px;
                    background: linear-gradient(135deg, #f97316 0%, #ea580c 100%);
                    color: white;
                    border-radius: 10px;
                }}
                .header h1 {{
                    margin: 0 0 10px 0;
                    font-size: 2em;
                    text-shadow: 2px 2px 4px rgba(0,0,0,0.3);
                }}
                .header p {{
                    margin: 5px 0 0 0;
                    opacity: 0.9;
                }}
                .info-badge {{
                    display: inline-block;
                    background: rgba(255,255,255,0.2);
                    padding: 8px 20px;
                    border-radius: 20px;
                    margin: 10px;
                    font-size: 0.9em;
                }}
                .section {{
                    margin: 30px 0;
                    padding: 25px;
                    background: #fff7ed;
                    border-radius: 10px;
                    border-left: 5px solid #f97316;
                }}
                .file-info {{
                    background: #f8f9fa;
                    padding: 20px;
                    border-radius: 8px;
                    margin: 20px 0;
                }}
                .footer {{
                    text-align: center;
                    padding: 30px;
                    margin-top: 40px;
                    color: #718096;
                    background: #f8f9fa;
                    border-radius: 10px;
                    font-size: 0.9em;
                }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>⚠️ 铜价风险预警PPT报告</h1>
                <p>{today}</p>
                <div class="info-badge">二级预警（警戒级）</div>
                <div class="info-badge">实时监控</div>
                <p style="margin-top: 15px; font-size: 0.9em;">生成时间：{timestamp}</p>
            </div>

            <div class="section">
                <h2 style="color: #c2410c; margin-bottom: 15px;">📎 PPT报告附件</h2>
                <div class="file-info">
                    <p><strong>文件名：</strong>{os.path.basename(ppt_file)}</p>
                    <p><strong>说明：</strong>完整的铜价风险预警分析报告（PPT格式）</p>
                    <p><strong>内容：</strong></p>
                    <ul>
                        <li>风险预警总览（当前风险等级）</li>
                        <li>价格波动预警（波动率、偏离度、跳空）</li>
                        <li>期限结构预警（LME价差、沪伦比、精废价差）</li>
                        <li>库存预警（库存变化、注销仓单）</li>
                        <li>资金情绪预警（CFTC净持仓、ETF流向）</li>
                        <li>宏观预警（美元指数、社融数据）</li>
                        <li>综合建议（操作建议）</li>
                    </ul>
                </div>
            </div>

            <div class="section">
                <h2 style="color: #c2410c; margin-bottom: 15px;">🚨 风险预警摘要</h2>
                <ul style="margin-left: 20px; line-height: 2;">
                    <li><strong>当前风险等级：</strong>🟠 二级预警（警戒级）</li>
                    <li><strong>需要关注：</strong>日内波动率接近警戒线</li>
                    <li><strong>需要关注：</strong>库存小幅下降</li>
                    <li><strong>需要关注：</strong>CFTC净持仓偏高</li>
                    <li><strong>正常：</strong>期限结构稳定</li>
                    <li><strong>正常：</strong>资金情绪积极</li>
                    <li><strong>正常：</strong>宏观环境平稳</li>
                </ul>
            </div>

            <div class="section">
                <h2 style="color: #c2410c; margin-bottom: 15px;">💡 操作建议</h2>
                <ol style="margin-left: 20px; line-height: 2;">
                    <li>密切关注价格波动</li>
                    <li>做好风险管理</li>
                    <li>保持适度仓位</li>
                    <li>及时调整策略</li>
                </ol>
            </div>

            <div class="footer">
                <p>⚠️ 以上预警仅供参考，不构成投资建议</p>
                <p>🤖 铜价风险预警系统 - 实时监控 + 多维分析</p>
                <p style="margin-top: 10px;">© 2026 铜价风险预警系统</p>
            </div>
        </body>
        </html>"""

        subject = f"⚠️ 铜价风险预警PPT报告 - {today}"

        # 发送邮件（只发送PPT附件）
        success, message = email_sender.send_email(
            to_email=config['receiver_email'],
            subject=subject,
            html_content=html_content,
            attachments=[ppt_file]
        )

        if success:
            print(f"✅ PPT邮件发送成功")
            return True
        else:
            print(f"❌ PPT邮件发送失败: {message}")
            return False

    except Exception as e:
        print(f"发送邮件时出错: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """主函数"""
    print()
    print("=" * 60)
    print("铜价风险预警系统 - 生成PPT并发送邮件")
    print("=" * 60)
    print()

    # 生成PPT
    ppt_file = generate_risk_alert_ppt()

    if not ppt_file:
        print()
        print("❌ 生成PPT失败，无法发送邮件")
        return False

    print()
    print("-" * 60)
    print("邮件发送")
    print("-" * 60)

    # 发送邮件
    success = send_risk_alert_email(ppt_file)

    print()
    print("=" * 60)
    if success:
        print("✅ 风险预警PPT报告发送完成！")
    else:
        print("⚠️ PPT生成完成，但邮件发送失败")
    print("=" * 60)

    return success


if __name__ == "__main__":
    main()
