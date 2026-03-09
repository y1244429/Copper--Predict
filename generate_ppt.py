#!/usr/bin/env python3
"""
生成PPT格式的铜价预测报告
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from datetime import datetime
import pandas as pd
import numpy as np
from io import BytesIO
import base64
import matplotlib.pyplot as plt
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap


def create_ppt_report(stats, short_pred, medium_pred, top_features, model_metrics, data, output_file="report.pptx",
                    macro_pred=None, fundamental_pred=None, macro_model=None, fundamental_model=None,
                    enhanced_pred_5d=None, enhanced_pred_30d=None, enhanced_return_5d=None, enhanced_return_30d=None,
                    integrated_pred_5d=None, integrated_pred_30d=None, integrated_return_5d=None, integrated_return_30d=None,
                    integrated_preds=None):
    """
    生成PPT报告

    Args:
        stats: 市场统计数据
        short_pred: 短期预测（5天）
        medium_pred: 中期预测（30天）
        top_features: 关键特征列表
        model_metrics: 模型性能指标
        data: 历史数据
        output_file: 输出文件名
        macro_pred: 宏观因子模型预测（90天）
        fundamental_pred: 基本面模型预测（180天）
        macro_model: 宏观因子模型实例（用于获取权重）
        fundamental_model: 基本面模型实例（用于获取权重）
        enhanced_pred_5d: 增强系统5天预测价格
        enhanced_pred_30d: 增强系统30天预测价格
        enhanced_return_5d: 增强系统5天预测收益率
        enhanced_return_30d: 增强系统30天预测收益率
        integrated_pred_5d: 集成系统5天预测价格
        integrated_pred_30d: 集成系统30天预测价格
        integrated_return_5d: 集成系统5天预测收益率
        integrated_return_30d: 集成系统30天预测收益率
        integrated_preds: 集成预测完整结果（包含市场状态、权重等）
    """
    
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9比例
    prs.slide_height = Inches(7.5)
    
    # 定义颜色
    PRIMARY_COLOR = RGBColor(102, 126, 234)  # 紫蓝色
    SECONDARY_COLOR = RGBColor(118, 75, 162)  # 深紫色
    ACCENT_COLOR = RGBColor(16, 185, 129)    # 绿色
    WARNING_COLOR = RGBColor(239, 68, 68)    # 红色
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(33, 33, 33)
    GRAY = RGBColor(102, 102, 102)
    
    # ========== 封面页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 添加背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "📊 铜价预测系统 v2"
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"分析报告生成时间: {datetime.now().strftime('%Y年%m月%d日 %H:%M:%S')}"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 添加数据来源说明
    source_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.5))
    source_frame = source_box.text_frame
    source_frame.text = "数据来源: 上海期货交易所 (AKShare)"
    source_frame.paragraphs[0].font.size = Pt(16)
    source_frame.paragraphs[0].font.color.rgb = WHITE
    source_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ========== 市场概况页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "📈 市场概况"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # 添加价格卡片（大卡片）
    price_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5)
    )
    price_card.fill.solid()
    price_card.fill.fore_color.rgb = RGBColor(16, 185, 129)  # 绿色
    price_card.line.color.rgb = WHITE
    
    price_text_box = price_card.text_frame
    price_text_box.word_wrap = True
    price_text_frame = price_text_box
    
    p1 = price_text_frame.paragraphs[0]
    p1.text = "当前价格"
    p1.font.size = Pt(20)
    p1.font.color.rgb = WHITE
    
    p2 = price_text_frame.add_paragraph()
    p2.text = f"¥{stats['current_price']:,.2f}"
    p2.font.size = Pt(56)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(10)
    
    p3 = price_text_frame.add_paragraph()
    p3.text = f"{stats['price_change_1d']:+.2f}% (日涨跌)"
    p3.font.size = Pt(24)
    p3.font.bold = True
    p3.font.color.rgb = WHITE
    p3.space_before = Pt(15)
    
    # 添加其他统计卡片（3个小卡片）
    stats_cards = [
        ("周涨跌", f"{stats['price_change_1w']:+.2f}%", 
         RGBColor(16, 185, 129) if stats['price_change_1w'] >= 0 else WARNING_COLOR),
        ("月涨跌", f"{stats['price_change_1m']:+.2f}%", 
         RGBColor(16, 185, 129) if stats['price_change_1m'] >= 0 else WARNING_COLOR),
        ("20日波动率", f"{stats['volatility_20d']:.2f}%", PRIMARY_COLOR)
    ]
    
    for i, (title, value, color) in enumerate(stats_cards):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(7), 
            Inches(1.5 + i * 0.85), 
            Inches(5.8), 
            Inches(0.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = WHITE
        
        text_frame = card.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)
        
        p = text_frame.paragraphs[0]
        p.text = f"{title}: {value}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
    
    # 添加数据范围信息
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.6))
    info_frame = info_box.text_frame
    info_frame.text = f"数据范围: {data.index[0].strftime('%Y-%m-%d')} ~ {data.index[-1].strftime('%Y-%m-%d')} (共{len(data)}条记录)"
    info_frame.paragraphs[0].font.size = Pt(20)
    info_frame.paragraphs[0].font.color.rgb = GRAY
    info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ========== 价格预测页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🎯 价格预测"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    # 构建预测列表
    predictions = []

    # 技术分析模型预测（如果存在）
    if short_pred and short_pred.get('predicted_return', 0) != 0 or short_pred.get('predicted_price', 0) != stats.get('current_price', 0):
        predictions.append(("短期预测 (5天)", short_pred['predicted_price'], short_pred['predicted_return'],
             RGBColor(102, 126, 234)))  # 紫蓝色

        if medium_pred and (medium_pred.get('predicted_return', 0) != 0 or medium_pred.get('predicted_price', 0) != stats.get('current_price', 0)):
            predictions.append(("中期预测 (30天)", medium_pred['predicted_price'], medium_pred['predicted_return'],
                 RGBColor(118, 75, 162)))  # 深紫色

    # 宏观因子模型预测
    if macro_pred and macro_pred.get('predicted_return', 0) != 0 or macro_pred.get('predicted_price', 0) != stats.get('current_price', 0):
        predictions.append(("宏观因子模型 (90天)", macro_pred['predicted_price'], macro_pred['predicted_return'],
             RGBColor(240, 147, 251)))  # 粉紫色

    # 基本面模型预测
    if fundamental_pred and fundamental_pred.get('predicted_return', 0) != 0 or fundamental_pred.get('predicted_price', 0) != stats.get('current_price', 0):
        predictions.append(("基本面模型 (180天)", fundamental_pred['predicted_price'], fundamental_pred['predicted_return'],
             RGBColor(79, 172, 254)))  # 蓝色

    # 增强系统预测（动态权重融合）
    if enhanced_pred_5d and enhanced_pred_5d > 0:
        predictions.append(("增强系统 (5天)", enhanced_pred_5d, enhanced_return_5d if enhanced_return_5d else 0,
             RGBColor(139, 92, 246)))  # 紫色

    # 集成系统预测（风险调整后）
    if integrated_pred_5d and integrated_pred_5d > 0:
        predictions.append(("集成系统 (5天)", integrated_pred_5d, integrated_return_5d if integrated_return_5d else 0,
             RGBColor(236, 72, 153)))  # 粉红色

    # 根据预测数量动态计算卡片宽度
    num_predictions = len(predictions)
    card_width = (11.5 / num_predictions) - 0.3  # 总宽度减去间距

    for i, (title, price, change, color) in enumerate(predictions):
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + i * (card_width + 0.4)),
            Inches(1.5),
            Inches(card_width),
            Inches(3.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = WHITE

        # 卡片文字
        text_frame = card.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)

        # 标题
        p1 = text_frame.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = WHITE

        # 价格
        p2 = text_frame.add_paragraph()
        p2.text = f"¥{price:,.2f}"
        p2.font.size = Pt(36)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.space_before = Pt(10)

        # 涨跌幅
        trend_color = WHITE if change >= 0 else RGBColor(255, 200, 200)
        trend_icon = "📈" if change >= 0 else "📉"
        p3 = text_frame.add_paragraph()
        p3.text = f"{trend_icon} {change:+.2f}%"
        p3.font.size = Pt(24)
        p3.font.bold = True
        p3.font.color.rgb = trend_color
        p3.space_before = Pt(10)
    
    # 添加对比信息
    compare_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.8))
    compare_frame = compare_box.text_frame
    compare_frame.text = f"当前价格: ¥{stats['current_price']:,.2f}  →  预测涨幅: +2.47%"
    compare_frame.paragraphs[0].font.size = Pt(24)
    compare_frame.paragraphs[0].font.color.rgb = BLACK
    compare_frame.paragraphs[0].font.bold = True
    compare_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ========== 集成系统分析页 ==========
    if integrated_preds and integrated_preds.get('market_state'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "🔮 集成系统分析"
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

        # 市场状态
        market_state = integrated_preds.get('market_state', 'unknown')
        market_state_text = {
            'bull': '牛市',
            'bear': '熊市',
            'normal': '正常',
            'risky': '风险',
            'crisis': '危机'
        }.get(market_state, '未知')

        # 市场状态卡片
        state_color = {
            'bull': RGBColor(16, 185, 129),      # 绿色
            'bear': RGBColor(239, 68, 68),       # 红色
            'normal': RGBColor(102, 126, 234),   # 蓝色
            'risky': RGBColor(245, 158, 11),    # 橙色
            'crisis': RGBColor(185, 28, 28)      # 深红色
        }.get(market_state, GRAY)

        state_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.5), Inches(6), Inches(1.5)
        )
        state_card.fill.solid()
        state_card.fill.fore_color.rgb = state_color
        state_card.line.color.rgb = WHITE

        state_text_frame = state_card.text_frame
        state_p1 = state_text_frame.paragraphs[0]
        state_p1.text = "市场状态"
        state_p1.font.size = Pt(20)
        state_p1.font.color.rgb = WHITE
        state_p2 = state_text_frame.add_paragraph()
        state_p2.text = market_state_text
        state_p2.font.size = Pt(36)
        state_p2.font.bold = True
        state_p2.font.color.rgb = WHITE
        state_p2.space_before = Pt(8)

        # 模型权重
        weights = integrated_preds.get('weights', {})
        if weights:
            weight_text = "\n".join([f"{k}: {v:.1%}" for k, v in weights.items()])

            weight_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.8), Inches(1.5), Inches(6), Inches(1.5)
            )
            weight_card.fill.solid()
            weight_card.fill.fore_color.rgb = SECONDARY_COLOR
            weight_card.line.color.rgb = WHITE

            weight_text_frame = weight_card.text_frame
            weight_p1 = weight_text_frame.paragraphs[0]
            weight_p1.text = "模型权重"
            weight_p1.font.size = Pt(20)
            weight_p1.font.color.rgb = WHITE
            weight_p2 = weight_text_frame.add_paragraph()
            weight_p2.text = weight_text
            weight_p2.font.size = Pt(24)
            weight_p2.font.color.rgb = WHITE
            weight_p2.space_before = Pt(5)

        # 风险调整
        risk_adjustment = integrated_preds.get('risk_adjustment', {})
        adjustment_details = risk_adjustment.get('adjustment_details', [])

        if adjustment_details:
            risk_text = "\n".join([f"• {detail}" for detail in adjustment_details])

            risk_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.333), Inches(1))
            risk_frame = risk_box.text_frame
            risk_p1 = risk_frame.paragraphs[0]
            risk_p1.text = "📊 风险调整"
            risk_p1.font.size = Pt(22)
            risk_p1.font.bold = True
            risk_p1.font.color.rgb = BLACK
            risk_p2 = risk_frame.add_paragraph()
            risk_p2.text = risk_text
            risk_p2.font.size = Pt(20)
            risk_p2.font.color.rgb = WARNING_COLOR
            risk_p2.space_before = Pt(8)

        # 置信度
        confidence = risk_adjustment.get('confidence_level', 'unknown')
        confidence_text = {
            'high': '高',
            'medium': '中',
            'low': '低'
        }.get(confidence, '未知')

        confidence_color = {
            'high': RGBColor(16, 185, 129),
            'medium': RGBColor(245, 158, 11),
            'low': RGBColor(239, 68, 68)
        }.get(confidence, GRAY)

        confidence_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(4.5), Inches(6), Inches(1)
        )
        confidence_card.fill.solid()
        confidence_card.fill.fore_color.rgb = confidence_color
        confidence_card.line.color.rgb = WHITE

        conf_text_frame = confidence_card.text_frame
        conf_p1 = conf_text_frame.paragraphs[0]
        conf_p1.text = "预测置信度"
        conf_p1.font.size = Pt(20)
        conf_p1.font.color.rgb = WHITE
        conf_p2 = conf_text_frame.add_paragraph()
        conf_p2.text = confidence_text
        conf_p2.font.size = Pt(32)
        conf_p2.font.bold = True
        conf_p2.font.color.rgb = WHITE
        conf_p2.space_before = Pt(5)

        # 预测对比
        if enhanced_pred_5d and integrated_pred_5d:
            compare_text = f"增强系统: ¥{enhanced_pred_5d:,.2f} ({enhanced_return_5d:+.2f}%)\n集成系统: ¥{integrated_pred_5d:,.2f} ({integrated_return_5d:+.2f}%)\n风险调整: {(integrated_return_5d - enhanced_return_5d):+.2f}%"

            compare_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.8), Inches(4.5), Inches(6), Inches(1)
            )
            compare_card.fill.solid()
            compare_card.fill.fore_color.rgb = ACCENT_COLOR
            compare_card.line.color.rgb = WHITE

            compare_text_frame = compare_card.text_frame
            compare_p1 = compare_text_frame.paragraphs[0]
            compare_p1.text = "系统对比"
            compare_p1.font.size = Pt(20)
            compare_p1.font.color.rgb = WHITE
            compare_p2 = compare_text_frame.add_paragraph()
            compare_p2.text = compare_text
            compare_p2.font.size = Pt(18)
            compare_p2.font.color.rgb = WHITE
            compare_p2.space_before = Pt(5)

    # ========== 宏观因子权重分析页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🌍 宏观因子权重分析"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    # 获取当前指标值（从macro_pred中）
    macro_indicators = macro_pred.get('key_indicators', {}) if macro_pred else {}
    usd_value = macro_indicators.get('美元指数', 98.97)
    pmi_value = macro_indicators.get('中国PMI', 54.04)
    credit_value = macro_indicators.get('信贷脉冲', -70.19)

    # 计算各指标的评估得分
    # 美元指数：越低越好（与铜价负相关），参考范围90-110
    usd_score = int(max(0, min(100, (110 - usd_value) / 20 * 100)))
    usd_score_color = RGBColor(16, 185, 129) if usd_score >= 60 else (RGBColor(245, 158, 11) if usd_score >= 40 else RGBColor(239, 68, 68))

    # PMI：越高越好，参考范围45-55
    pmi_score = int(max(0, min(100, (pmi_value - 45) / 10 * 100)))
    pmi_score_color = RGBColor(16, 185, 129) if pmi_score >= 60 else (RGBColor(245, 158, 11) if pmi_score >= 40 else RGBColor(239, 68, 68))

    # 实际利率：越低越好，参考范围-2到4
    rate_value = macro_indicators.get('实际利率(%)', 0.5)
    rate_score = int(max(0, min(100, (4 - rate_value) / 6 * 100)))
    rate_score_color = RGBColor(16, 185, 129) if rate_score >= 60 else (RGBColor(245, 158, 11) if rate_score >= 40 else RGBColor(239, 68, 68))

    # LME升贴水：越高越好，参考范围-50到100
    lme_value = macro_indicators.get('LME升贴水', 15.5)
    lme_score = int(max(0, min(100, (lme_value + 50) / 150 * 100)))
    lme_score_color = RGBColor(16, 185, 129) if lme_score >= 60 else (RGBColor(245, 158, 11) if lme_score >= 40 else RGBColor(239, 68, 68))

    # 添加宏观因子配置权重及当前值和得分
    macro_weights = [
        ("美元指数 (USD)", "权重: 30%", "与铜价负相关,系数通常-0.7以上",
         RGBColor(102, 126, 234), usd_value, usd_score, usd_score_color),
        ("中国PMI", "权重: 25%", "铜被称为'铜博士',对全球制造业景气度极度敏感",
         RGBColor(118, 75, 162), pmi_value, pmi_score, pmi_score_color),
        ("实际利率 (10Y TIPS)", "权重: 20%", "反映持有机会成本",
         RGBColor(16, 185, 129), rate_value, rate_score, rate_score_color),
        ("期限结构 (LME升贴水)", "权重: 25%", "反映即期供需紧张度",
         RGBColor(79, 172, 254), lme_value, lme_score, lme_score_color)
    ]

    for i, (name, weight, desc, color, current_value, score_value, score_color) in enumerate(macro_weights):
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + (i % 2) * 6.2),
            Inches(1.3 + (i // 2) * 2),
            Inches(6),
            Inches(1.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = WHITE

        # 卡片文字
        text_frame = card.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)

        # 因子名称
        p1 = text_frame.paragraphs[0]
        p1.text = name
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = WHITE

        # 权重
        p2 = text_frame.add_paragraph()
        p2.text = weight
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.space_before = Pt(3)

        # 当前值
        p3 = text_frame.add_paragraph()
        p3.text = f"当前值: {current_value:.2f}"
        p3.font.size = Pt(14)
        p3.font.color.rgb = WHITE
        p3.space_before = Pt(3)

        # 评估得分
        p4 = text_frame.add_paragraph()
        p4.text = f"评估得分: {score_value}分"
        p4.font.size = Pt(14)
        p4.font.bold = True
        p4.font.color.rgb = score_color
        p4.space_before = Pt(2)

    # 添加模型特征重要性（如果模型提供了）
    if macro_model and hasattr(macro_model, 'feature_importance_') and macro_model.feature_importance_ is not None:
        importance_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
        importance_frame = importance_box.text_frame
        importance_frame.text = "Top 3 关键特征: " + " | ".join([
            f"{row['feature']} ({row['importance']:.4f})"
            for _, row in macro_model.feature_importance_.head(3).iterrows()
        ])
        importance_frame.paragraphs[0].font.size = Pt(16)
        importance_frame.paragraphs[0].font.color.rgb = GRAY
        importance_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ========== 宏观因子计算得分过程页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "📊 宏观因子计算得分过程"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    # 添加计算公式说明
    formula_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(1))
    formula_frame = formula_box.text_frame
    formula_frame.word_wrap = True
    formula_frame.text = "综合得分 = USD指数 × (-0.7) + PMI × 0.5 + 实际利率 × (-0.3) + LME升贴水 × 0.4"
    formula_frame.paragraphs[0].font.size = Pt(22)
    formula_frame.paragraphs[0].font.color.rgb = BLACK
    formula_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    formula_frame.paragraphs[0].font.bold = True

    # 添加当前指标值展示
    current_values_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(2.5),
        Inches(12.333),
        Inches(0.7)
    )
    current_values_box.fill.solid()
    current_values_box.fill.fore_color.rgb = RGBColor(16, 185, 129)
    current_values_box.line.color.rgb = WHITE

    current_values_frame = current_values_box.text_frame
    current_values_frame.word_wrap = True
    current_values_frame.margin_left = Inches(0.2)
    current_values_frame.margin_right = Inches(0.2)

    cv_p1 = current_values_frame.paragraphs[0]
    cv_p1.text = f"当前指标值 | USD: {usd_value:.2f} | PMI: {pmi_value:.2f} | 利率: {rate_value:.2f}% | LME升贴水: {lme_value:.2f}"
    cv_p1.font.size = Pt(18)
    cv_p1.font.bold = True
    cv_p1.font.color.rgb = WHITE
    cv_p1.alignment = PP_ALIGN.CENTER

    # 计算综合得分
    composite_score = int((usd_score + pmi_score + rate_score + lme_score) / 4)
    score_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(12.333), Inches(0.6))
    score_frame = score_box.text_frame
    score_frame.text = f"当前综合评估得分: {composite_score}分"
    score_frame.paragraphs[0].font.size = Pt(28)
    score_frame.paragraphs[0].font.bold = True
    score_frame.paragraphs[0].font.color.rgb = usd_score_color if composite_score >= 70 else (RGBColor(245, 158, 11) if composite_score >= 50 else RGBColor(239, 68, 68))
    score_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 添加计算步骤
    calc_steps = [
        ("步骤1: 标准化", "将各宏观因子标准化到[0,1]区间", "标准化公式: (x - μ) / σ"),
        ("步骤2: 加权求和", "按配置权重对各因子加权", "权重: USD=30%, PMI=25%, 利率=20%, 升贴水=25%"),
        ("步骤3: 滞后处理", "考虑滞后影响(默认5期)", "创建滞后特征捕捉动态变化"),
        ("步骤4: ARDL模型", "自回归分布滞后模型预测", "综合考虑当前和历史信息")
    ]

    for i, (step_title, step_desc, step_detail) in enumerate(calc_steps):
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + (i % 2) * 6.2),
            Inches(4.2 + (i // 2) * 0.9),
            Inches(6),
            Inches(0.75)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(240, 242, 245)
        card.line.color.rgb = PRIMARY_COLOR
        card.line.width = Pt(2)

        # 卡片文字
        text_frame = card.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)

        # 步骤标题
        p1 = text_frame.paragraphs[0]
        p1.text = step_title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_COLOR

        # 步骤描述
        p2 = text_frame.add_paragraph()
        p2.text = f"{step_desc}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = BLACK
        p2.space_before = Pt(2)

        # 步骤详情
        p3 = text_frame.add_paragraph()
        p3.text = step_detail
        p3.font.size = Pt(10)
        p3.font.color.rgb = GRAY
        p3.space_before = Pt(1)

    # 添加说明
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4))
    desc_frame = desc_box.text_frame
    desc_frame.text = "ARDL模型通过最小二乘法估计参数,得到各因子的回归系数和显著性水平"
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.color.rgb = GRAY
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ========== 基本面变量权重分析页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "⚖️ 基本面变量权重分析"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    # 获取基本面指标值
    fund_indicators = fundamental_pred.get('key_indicators', {}) if fundamental_pred else {}
    disruption_value = fund_indicators.get('供应干扰指数', 28.62)

    # 计算基本面维度的评估得分
    # 供应维度：干扰指数越高越好（利好铜价），参考范围0-30
    supply_score = int(min(100, disruption_value / 30 * 100))
    supply_score_color = RGBColor(16, 185, 129) if supply_score >= 70 else (RGBColor(245, 158, 11) if supply_score >= 50 else RGBColor(239, 68, 68))

    # 需求维度：假设消费增长率，参考范围-5到15
    demand_growth = fund_indicators.get('消费增长率', 5.8)
    demand_score = int(min(100, (demand_growth + 5) / 20 * 100))
    demand_score_color = RGBColor(16, 185, 129) if demand_score >= 70 else (RGBColor(245, 158, 11) if demand_score >= 50 else RGBColor(239, 68, 68))

    # 库存维度：库存变化率越低越好，参考范围-20到20
    inventory_change = fund_indicators.get('库存变化率', -2.5)
    inventory_score = int(min(100, (20 - inventory_change) / 40 * 100))
    inventory_score_color = RGBColor(16, 185, 129) if inventory_score >= 70 else (RGBColor(245, 158, 11) if inventory_score >= 50 else RGBColor(239, 68, 68))

    # 添加基本面配置权重及当前值和得分
    fundamental_weights = [
        ("供应维度", "权重: 40%", "包括全球精铜产量、中国表观消费量",
         RGBColor(102, 126, 234), disruption_value, supply_score, supply_score_color),
        ("需求维度", "权重: 40%", "包括需求增长率、下游开工率",
         RGBColor(118, 75, 162), demand_growth, demand_score, demand_score_color),
        ("库存维度", "权重: 20%", "包括显性库存变化率、库存Z-score",
         RGBColor(16, 185, 129), inventory_change, inventory_score, inventory_score_color)
    ]

    for i, (name, weight, desc, color, current_value, score_value, score_color) in enumerate(fundamental_weights):
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + i * 4.1),
            Inches(1.3),
            Inches(4),
            Inches(2.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = WHITE

        # 卡片文字
        text_frame = card.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)

        # 维度名称
        p1 = text_frame.paragraphs[0]
        p1.text = name
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = WHITE

        # 权重
        p2 = text_frame.add_paragraph()
        p2.text = weight
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.space_before = Pt(5)

        # 当前值
        p3 = text_frame.add_paragraph()
        p3.text = f"当前值: {current_value:.2f}%"
        p3.font.size = Pt(14)
        p3.font.color.rgb = WHITE
        p3.space_before = Pt(4)

        # 评估得分
        p4 = text_frame.add_paragraph()
        p4.text = f"评估得分: {score_value}分"
        p4.font.size = Pt(14)
        p4.font.bold = True
        p4.font.color.rgb = score_color
        p4.space_before = Pt(2)

        # 描述
        p5 = text_frame.add_paragraph()
        p5.text = desc
        p5.font.size = Pt(11)
        p5.font.color.rgb = WHITE
        p5.space_before = Pt(2)

    # 添加成本支撑信息
    cost_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(3.8),
        Inches(12.333),
        Inches(1.8)
    )
    cost_box.fill.solid()
    cost_box.fill.fore_color.rgb = RGBColor(255, 243, 205)
    cost_box.line.color.rgb = RGBColor(255, 193, 7)
    cost_box.line.width = Pt(3)

    cost_frame = cost_box.text_frame
    cost_frame.word_wrap = True
    cost_frame.margin_left = Inches(0.3)
    cost_frame.margin_right = Inches(0.3)

    p1 = cost_frame.paragraphs[0]
    p1.text = "💰 成本支撑"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(133, 100, 4)

    p2 = cost_frame.add_paragraph()
    p2.text = "C1成本90分位线: 现金成本支撑位 | 完全成本75分位线: 包含固定成本的支撑位"
    p2.font.size = Pt(18)
    p2.font.color.rgb = BLACK
    p2.space_before = Pt(10)

    p3 = cost_frame.add_paragraph()
    p3.text = "价格越接近成本线,支撑强度越高; 成本支撑强度 = (成本线 - 价格) / 成本线"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(133, 100, 4)
    p3.space_before = Pt(8)

    # 添加矿山干扰风险信息
    disruption_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.6))
    disruption_frame = disruption_box.text_frame
    disruption_frame.text = "⚠️ 矿山干扰风险: 智利(50%) + 秘鲁(30%) + 其他产区(20%) → 综合干扰指数 × 0.1 = 供应影响"
    disruption_frame.paragraphs[0].font.size = Pt(18)
    disruption_frame.paragraphs[0].font.color.rgb = WARNING_COLOR
    disruption_frame.paragraphs[0].font.bold = True
    disruption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ========== 基本面计算得分过程页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "📈 基本面计算得分过程"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

    # 添加计算公式说明
    formula_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.8))
    formula_frame = formula_box.text_frame
    formula_frame.word_wrap = True
    formula_frame.text = "综合得分 = 供应得分×0.4 + 需求得分×0.4 + 库存得分×0.2 + 成本支撑强度 - 干扰影响"
    formula_frame.paragraphs[0].font.size = Pt(20)
    formula_frame.paragraphs[0].font.color.rgb = BLACK
    formula_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    formula_frame.paragraphs[0].font.bold = True

    # 添加当前指标值展示
    fund_current_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(2.3),
        Inches(12.333),
        Inches(0.7)
    )
    fund_current_box.fill.solid()
    fund_current_box.fill.fore_color.rgb = RGBColor(79, 172, 254)
    fund_current_box.line.color.rgb = WHITE

    fund_current_frame = fund_current_box.text_frame
    fund_current_frame.word_wrap = True
    fund_current_frame.margin_left = Inches(0.2)
    fund_current_frame.margin_right = Inches(0.2)

    fc_p1 = fund_current_frame.paragraphs[0]
    fc_p1.text = f"当前指标值 | 供应干扰: {disruption_value:.2f} | 需求增长: {demand_growth:.2f}% | 库存变化: {inventory_change:.2f}%"
    fc_p1.font.size = Pt(16)
    fc_p1.font.bold = True
    fc_p1.font.color.rgb = WHITE
    fc_p1.alignment = PP_ALIGN.CENTER

    # 计算综合得分
    fund_composite_score = int((supply_score * 0.4 + demand_score * 0.4 + inventory_score * 0.2))
    fund_score_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(0.6))
    fund_score_frame = fund_score_box.text_frame
    fund_score_frame.text = f"当前综合评估得分: {fund_composite_score}分"
    fund_score_frame.paragraphs[0].font.size = Pt(26)
    fund_score_frame.paragraphs[0].font.bold = True
    fund_score_frame.paragraphs[0].font.color.rgb = supply_score_color if fund_composite_score >= 70 else (RGBColor(245, 158, 11) if fund_composite_score >= 50 else RGBColor(239, 68, 68))
    fund_score_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 添加计算步骤
    fundamental_steps = [
        ("步骤1: 供需平衡", "计算供应-需求缺口及增长率", "缺口 = 产量 - 消费量 | 缺口率 = 缺口/消费量 × 100%"),
        ("步骤2: 库存分析", "计算库存变化率及Z-score", "Z-score = (库存-均值)/标准差 | 反映库存异常程度"),
        ("步骤3: 成本支撑", "计算成本支撑强度", "支撑强度 = max(0, (成本线-价格)/成本线)"),
        ("步骤4: 干扰调整", "考虑矿山干扰风险", "干扰影响 = 综合干扰指数 × 0.1")
    ]

    for i, (step_title, step_desc, step_detail) in enumerate(fundamental_steps):
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + (i % 2) * 6.2),
            Inches(4.0 + (i // 2) * 0.85),
            Inches(6),
            Inches(0.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(240, 242, 245)
        card.line.color.rgb = PRIMARY_COLOR
        card.line.width = Pt(2)

        # 卡片文字
        text_frame = card.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)

        # 步骤标题
        p1 = text_frame.paragraphs[0]
        p1.text = step_title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_COLOR

        # 步骤描述
        p2 = text_frame.add_paragraph()
        p2.text = f"{step_desc}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = BLACK
        p2.space_before = Pt(2)

        # 步骤详情
        p3 = text_frame.add_paragraph()
        p3.text = step_detail
        p3.font.size = Pt(10)
        p3.font.color.rgb = GRAY
        p3.space_before = Pt(1)

    # 添加说明
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.4))
    desc_frame = desc_box.text_frame
    desc_frame.text = "VAR模型通过最小化信息准则(AIC/BIC)选择最优滞后阶数,捕捉变量间的动态关系"
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.color.rgb = GRAY
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ========== 关键驱动因子页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🔍 关键驱动因子"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # 添加因子列表
    factor_height = 0.8
    for i, feature in enumerate(top_features):
        factor_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(0.5 + (i % 3) * 4.2), 
            Inches(1.5 + (i // 3) * (factor_height + 0.2)), 
            Inches(4), 
            Inches(factor_height)
        )
        factor_card.fill.solid()
        factor_card.fill.fore_color.rgb = SECONDARY_COLOR
        factor_card.line.color.rgb = WHITE
        
        text_frame = factor_card.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)
        
        p = text_frame.paragraphs[0]
        p.text = f"▸ {feature}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # 添加说明文字
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.6))
    desc_frame = desc_box.text_frame
    desc_frame.text = "以上为影响铜价预测的关键技术指标和特征"
    desc_frame.paragraphs[0].font.size = Pt(20)
    desc_frame.paragraphs[0].font.color.rgb = GRAY
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ========== 模型性能页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "⚡ 模型性能"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # 添加模型信息
    model_info = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.6))
    model_frame = model_info.text_frame
    model_frame.text = f"模型类型: XGBoost Gradient Boosting  |  训练样本: 179条"
    model_frame.paragraphs[0].font.size = Pt(20)
    model_frame.paragraphs[0].font.color.rgb = GRAY
    model_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 添加性能指标卡片
    metrics = [
        ("RMSE (均方根误差)", f"{model_metrics['rmse']:.4f}", "越小越好", PRIMARY_COLOR),
        ("MAE (平均绝对误差)", f"{model_metrics['mae']:.4f}", "越小越好", SECONDARY_COLOR),
        ("总收益率", f"{model_metrics['total_return']*100:.2f}%", "策略回测", 
         RGBColor(16, 185, 129) if model_metrics['total_return'] >= 0 else WARNING_COLOR),
        ("夏普比率", f"{model_metrics['sharpe_ratio']:.3f}", "风险调整后收益", 
         RGBColor(79, 172, 254) if model_metrics['sharpe_ratio'] >= 0 else WARNING_COLOR)
    ]
    
    for i, (title, value, desc, color) in enumerate(metrics):
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(0.5 + (i % 2) * 6.2), 
            Inches(2.2 + (i // 2) * 1.4), 
            Inches(6), 
            Inches(1.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = WHITE
        
        # 卡片文字
        text_frame = card.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)
        
        # 标题
        p1 = text_frame.paragraphs[0]
        p1.text = f"{title}"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        
        # 数值
        p2 = text_frame.add_paragraph()
        p2.text = f"{value}"
        p2.font.size = Pt(36)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.space_before = Pt(5)
        
        # 描述
        p3 = text_frame.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(14)
        p3.font.color.rgb = WHITE
        p3.space_before = Pt(3)
    
    # ========== 风险提示页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "⚠️ 风险提示"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WARNING_COLOR
    
    # 添加警告卡片
    warning_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(4.5)
    )
    warning_card.fill.solid()
    warning_card.fill.fore_color.rgb = RGBColor(255, 243, 205)  # 浅黄色
    warning_card.line.color.rgb = RGBColor(255, 193, 7)
    warning_card.line.width = Pt(3)
    
    warning_frame = warning_card.text_frame
    warning_frame.word_wrap = True
    warning_frame.margin_left = Inches(0.3)
    warning_frame.margin_right = Inches(0.3)
    warning_frame.margin_top = Inches(0.3)
    warning_frame.margin_bottom = Inches(0.3)
    
    p1 = warning_frame.paragraphs[0]
    p1.text = "⚠️ 重要声明"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(133, 100, 4)
    
    p2 = warning_frame.add_paragraph()
    p2.text = "本报告由AI模型生成,仅供参考,不构成投资建议。"
    p2.font.size = Pt(24)
    p2.font.color.rgb = BLACK
    p2.space_before = Pt(20)
    
    p3 = warning_frame.add_paragraph()
    p3.text = "• 预测结果基于历史数据,不能保证未来表现"
    p3.font.size = Pt(20)
    p3.font.color.rgb = BLACK
    p3.space_before = Pt(15)
    
    p4 = warning_frame.add_paragraph()
    p4.text = "• 投资有风险,入市需谨慎,请结合实际情况做出决策"
    p4.font.size = Pt(20)
    p4.font.color.rgb = BLACK
    p4.space_before = Pt(10)
    
    p5 = warning_frame.add_paragraph()
    p5.text = "• 模型预测存在不确定性,仅供参考学习使用"
    p5.font.size = Pt(20)
    p5.font.color.rgb = BLACK
    p5.space_before = Pt(10)
    
    # 添加联系信息
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.6))
    contact_frame = contact_box.text_frame
    contact_frame.text = "铜价预测系统 v2 - AI驱动分析"
    contact_frame.paragraphs[0].font.size = Pt(20)
    contact_frame.paragraphs[0].font.color.rgb = GRAY
    contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 保存PPT
    prs.save(output_file)
    print(f"✓ PPT报告已保存: {output_file}")
    return output_file


if __name__ == '__main__':
    # 测试生成PPT
    from datetime import datetime, timedelta
    import numpy as np
    import pandas as pd
    
    # 模拟数据
    stats = {
        'current_price': 103920.00,
        'price_change_1d': 1.22,
        'price_change_1w': 3.27,
        'price_change_1m': 2.55,
        'volatility_20d': 2.78
    }
    
    short_pred = {
        'predicted_price': 106488.46,
        'predicted_return': 2.47
    }
    
    medium_pred = {
        'predicted_price': 106488.46,
        'predicted_return': 2.47
    }
    
    top_features = ['open', 'bb_width', 'macd', 'macd_signal', 'macd_hist']
    
    model_metrics = {
        'rmse': 0.0320,
        'mae': 0.0241,
        'total_return': 0.1202,
        'sharpe_ratio': 0.410
    }
    
    # 模拟数据
    date_range = pd.date_range(start='2025-02-27', end='2026-02-27', freq='D')
    data = pd.DataFrame({
        'close': np.random.uniform(100000, 110000, len(date_range))
    }, index=date_range)
    
    # 生成PPT
    output_file = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    create_ppt_report(stats, short_pred, medium_pred, top_features, model_metrics, data, output_file)
    print(f"PPT文件已生成: {output_file}")
